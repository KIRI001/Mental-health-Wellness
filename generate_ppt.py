import collections
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Creamy brown background: F5DEB3 (Wheat) or EEDC82 (Flax) or EADDCA (Almond)
BG_COLOR = RGBColor(234, 221, 202) # Creamy brown / Almond
FONT_NAME = 'Times New Roman'

prs = Presentation()

# Slide width and height (default 16:9)
# Slide formats typically use 16:9: 10 inches by 5.625 inches
# 1 inch = 914400 EMUs
# The font sizes are specified.

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_slide(title_text, content_text, layout_index=1):
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)
    
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = title_text
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.name = FONT_NAME
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER
            
    if layout_index == 1 and len(slide.shapes) > 1:
        body_shape = slide.shapes.placeholders[1]
        body_shape.text = content_text
        for paragraph in body_shape.text_frame.paragraphs:
            paragraph.font.name = FONT_NAME
            # If it looks like a heading
            if paragraph.text.startswith("- ") or paragraph.level > 0:
                paragraph.font.size = Pt(12)
            else:
                paragraph.font.size = Pt(14) # Maybe heading size for top level
            
            # Simple heuristic to apply sizes based on lines:
            
    return slide

def add_custom_slide(title_text, lines):
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)
    
    if slide.shapes.title:
        title_shape = slide.shapes.title
        title_shape.text = title_text
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.name = FONT_NAME
            paragraph.font.size = Pt(16)
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER

    # Add textbox
    left = Pt(50)
    top = Pt(100)
    width = prs.slide_width - Pt(100)
    height = prs.slide_height - Pt(150)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
            
        # Determine indentation/level
        if line.startswith("• "):
            p.level = 1
            line = line[2:]
        elif line.startswith("  - "):
            p.level = 2
            line = line[4:]
        
        # Check if heading
        if ":" in line and len(line) < 40 and not line.startswith("Accuracy"):
            # Might be a heading, wait, let's just make it bold if it has a colon at the end
            pass
        
        p.text = line
        p.font.name = FONT_NAME
        
        # Let's use 12pt for body, 14pt for headings
        if p.level == 0 and not line.startswith("Accuracy") and not line.startswith("Precision") and not line.startswith("Recall") and not line.startswith("F1") and not line.startswith("AUROC"):
            p.font.size = Pt(14)
            p.font.bold = True
        else:
            p.font.size = Pt(12)

# Slide 1
s1_lines = [
    "Student Name: [Your Name]",
    "Registration Number: [Your Reg No]",
    "Department of Computer Science & Engineering",
    "Dumka Engineering College",
    "Guide Name: [Guide Name]",
    "Academic Year: [Year]"
]
add_custom_slide("Automated Pityriasis Alba Detection Using Deep Learning", s1_lines)

# Slide 2
s2_lines = [
    "Pityriasis Alba is a common hypopigmented skin disorder affecting mainly children and adolescents.",
    "",
    "Importance of Early Diagnosis:",
    "• Essential for effective management and patient reassurance.",
    "",
    "Role of Artificial Intelligence:",
    "• AI enables automated, accurate, and rapid skin disease detection.",
    "• Reduces manual diagnostic errors and provides scalable healthcare solutions."
]
add_custom_slide("Introduction", s2_lines)

# Slide 3
s3_lines = [
    "Challenges in Current Diagnosis:",
    "• Difficulty in manual diagnosis due to subtle visual features",
    "• Similar appearance with Vitiligo, Eczema, and Tinea Versicolor",
    "• High dependency on experienced dermatologists",
    "• Limited accessibility to specialists in rural areas",
    "",
    "Solution:",
    "• A strong need for an automated, AI-driven diagnostic solution to assist medical professionals."
]
add_custom_slide("Problem Statement", s3_lines)

# Slide 4
s4_lines = [
    "Key Objectives:",
    "• Automated detection of Pityriasis Alba from skin lesion images",
    "• Comparison of deep learning architectures: CNN, DenseNet, and ResNet",
    "• Integration of the best model into a user-friendly web application",
    "• Improvement of diagnostic accessibility for remote locations",
    "• Comprehensive evaluation using multiple performance metrics"
]
add_custom_slide("Project Objectives", s4_lines)

# Slide 5
s5_lines = [
    "Dataset Details:",
    "• Total dataset size: 8,440 images",
    "",
    "Disease Categories:",
    "• Pityriasis Alba",
    "• Vitiligo",
    "• Eczema",
    "• Tinea Versicolor",
    "",
    "Significance:",
    "• High dataset diversity ensuring clinical relevance and model robustness"
]
add_custom_slide("Dataset Overview", s5_lines)

# Slide 6
s6_lines = [
    "Steps Taken:",
    "• Image resizing to standard dimensions for model input",
    "• Normalization of pixel values for faster convergence",
    "• Data augmentation to improve model generalization",
    "• Dataset balancing to handle class disparities",
    "• Training, validation, and testing split (e.g., 70:15:15)"
]
add_custom_slide("Data Preprocessing", s6_lines)

# Slide 7
s7_lines = [
    "Workflow:",
    "Image Upload → Preprocessing → Feature Extraction → Model Prediction → Disease Analysis → Result Generation",
    "",
    "[Placeholder for Visual Workflow Diagram]"
]
add_custom_slide("System Workflow", s7_lines)

# Slide 8
s8_lines = [
    "Architectures Explored:",
    "• CNN (Custom Convolutional Neural Network)",
    "• DenseNet",
    "• ResNet",
    "",
    "Why Multiple Architectures?",
    "• To compare lightweight custom models against deep pre-trained networks.",
    "• To identify the best balance between accuracy and computational efficiency for web deployment."
]
add_custom_slide("Deep Learning Models Used", s8_lines)

# Slide 9
s9_lines = [
    "CNN Components:",
    "• Convolution layers for localized feature extraction",
    "• Pooling layers for spatial dimension reduction",
    "• Fully connected layers for high-level reasoning",
    "• Binary/Multi-class classification output layer",
    "",
    "Highlight:",
    "• Excellent at capturing hierarchical spatial features from skin images."
]
add_custom_slide("CNN Architecture", s9_lines)

# Slide 10
s10_lines = [
    "DenseNet Features:",
    "• Dense connectivity where each layer receives inputs from all preceding layers",
    "• Feature reuse for maximized information flow",
    "• Efficient gradient flow mitigating the vanishing gradient problem",
    "• Improved learning efficiency with fewer parameters compared to traditional deep networks"
]
add_custom_slide("DenseNet Architecture", s10_lines)

# Slide 11
s11_lines = [
    "ResNet Features:",
    "• Residual learning framework",
    "• Skip connections that bypass some layers",
    "• Deep feature extraction without degradation",
    "• Improved optimization and training of very deep networks"
]
add_custom_slide("ResNet Architecture", s11_lines)

# Slide 12
s12_lines = [
    "Parameters:",
    "• Optimizer: Adam",
    "• Loss Function: Categorical / Binary Cross Entropy",
    "• Epochs: 10",
    "• Multiple independent runs for reliability",
    "• Tracked evaluation metrics across epochs to prevent overfitting"
]
add_custom_slide("Training Methodology", s12_lines)

# Slide 13
s13_lines = [
    "Metrics Defined:",
    "• Accuracy: Overall correctness of the model",
    "• Precision: Proportion of true positive results among all positive predictions",
    "• Recall: Ability to identify all actual positive cases",
    "• F1 Score: Harmonic mean of Precision and Recall",
    "• AUROC: Area Under the Receiver Operating Characteristic Curve",
    "",
    "Importance:",
    "• Critical in medical diagnosis to minimize false negatives (Recall) and false positives (Precision)."
]
add_custom_slide("Evaluation Metrics", s13_lines)

# Slide 14
s14_lines = [
    "Model Performance:",
    "",
    "CNN:",
    "• Accuracy: 78.84% | Precision: 71.50% | Recall: 95.31% | F1 Score: 81.66% | AUROC: 89.11%",
    "",
    "DenseNet:",
    "• Accuracy: 73.57% | Precision: 67.58% | Recall: 90.59% | F1 Score: 77.26% | AUROC: 81.89%",
    "",
    "ResNet:",
    "• Accuracy: 69.85% | Precision: 64.10% | Recall: 94.42% | F1 Score: 75.85% | AUROC: 85.83%",
    "",
    "Conclusion: CNN is the best-performing model."
]
add_custom_slide("Model Comparison Summary", s14_lines)

# Slide 15
s15_lines = [
    "[Placeholder for Confusion Matrix]",
    "[Placeholder for ROC Curve]",
    "",
    "Key Performance:",
    "• Achieved the highest overall accuracy (78.84%) and AUROC (89.11%).",
    "• Most balanced capability for Pityriasis Alba detection."
]
add_custom_slide("CNN Results", s15_lines)

# Slide 16
s16_lines = [
    "[Placeholder for Confusion Matrix]",
    "[Placeholder for ROC Curve]",
    "",
    "Key Performance:",
    "• Accuracy: 73.57%",
    "• Strong balanced classification capability."
]
add_custom_slide("DenseNet Results", s16_lines)

# Slide 17
s17_lines = [
    "[Placeholder for Confusion Matrix]",
    "[Placeholder for ROC Curve]",
    "",
    "Key Performance:",
    "• Accuracy: 69.85%",
    "• Strong disease detection capability through high recall (94.42%)."
]
add_custom_slide("ResNet Results", s17_lines)

# Slide 18
s18_lines = [
    "Key Platform Features:",
    "• Login Page & User Profile Page",
    "• Dashboard for Analytics",
    "• Image Upload Page",
    "• QR-based Mobile Upload",
    "• Prediction Result Page",
    "",
    "[Placeholder for Web Application Screenshots]"
]
add_custom_slide("Web Application", s18_lines)

# Slide 19
s19_lines = [
    "Conclusion:",
    "• Deep learning successfully detects Pityriasis Alba.",
    "• CNN achieved the best overall results (78.84% accuracy).",
    "• Practical deployment achieved through web integration.",
    "",
    "Future Scope:",
    "• Doctor consultation integration",
    "• Larger clinical datasets",
    "• Improved accuracy and robustness",
    "• Multi-disease classification",
    "• Mobile application deployment"
]
add_custom_slide("Conclusion and Future Scope", s19_lines)

# Slide 20
s20_lines = [
    "",
    "",
    "Questions and Discussion",
    "",
    "[Placeholder for Medical AI-themed Graphics]"
]
add_custom_slide("THANK YOU", s20_lines)

prs.save('/tmp/workspace/KIRI001/MHW-/Presentation.pptx')
print("Presentation generated successfully!")
